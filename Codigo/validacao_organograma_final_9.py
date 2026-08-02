# -----------------------------------------------------------------------------
# Validação de organograma (nome + vínculo + GESTOR) contra a base Fênix
# -----------------------------------------------------------------------------
# O gestor de cada caixa é lido do TEXTO ALTERNATIVO da forma no PowerPoint,
# no formato "Gestor=NOME DO GESTOR" (Botão direito na caixa -> Exibir Texto
# Alternativo). Caixas sem esse texto alternativo simplesmente não entram na
# validação de gestor (aparecem como "SEM GESTOR NO PPTX").
# -----------------------------------------------------------------------------
# Requisitos: pip install python-pptx pandas openpyxl xlrd
# -----------------------------------------------------------------------------

import re
import unicodedata
from difflib import SequenceMatcher
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pandas as pd

NS_P = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

# === CONFIGURAÇÃO: edite só estes três caminhos =============================
# IMPORTANTE: sempre use r"..." (raw string) em caminhos do Windows.
# Sem o "r", o Python tenta interpretar "\U", "\n", "\t" etc como caracteres
# especiais, e "...\Users\..." quebra com "SyntaxError: unicodeescape".
CAMINHO_PPTX = r"C:\Automacao\Organograma_Global Markets_Jul26.pptx"
CAMINHO_FENIX = r"C:\Automacao\Base_Fenix_Global_Markets_31.07.26.xlsx"
CAMINHO_SAIDA = r"C:\Automacao\relatorio_organograma_Global_Markets.xlsx"
# O topo de cada organograma normalmente é o único "Estatutário" (diretor) -
# e o PPT não escreve isso na caixa dele, só o cargo/área. Adicione aqui o(s)
# nome(s) do topo de CADA organograma que você validar (qualquer parte do
# nome, em maiúsculas, já resolve).
NOMES_ESTATUTARIO = ["DANIELA FASOLARI", "REGIS EDUARDO PRENHACA CARREIRA"]
# Colunas de abertura de área na base Fênix, da mais ampla pra mais específica
ABERTURA_COLS = ['Área\nGerencial', 'Abertura Área I', 'Abertura Área II',
                 'Abertura Área III', 'Abertura Área IV', 'Abertura Área V']
# =============================================================================


# ---------------------------------------------------------------------------
# Normalização e matching de nomes (usa a própria base Fênix como dicionário
# de nomes reais, para não depender de heurística de palavras-chave)
# ---------------------------------------------------------------------------

def strip_accents(s):
    return ''.join(c for c in unicodedata.normalize('NFD', s) if unicodedata.category(c) != 'Mn')

def norm(s):
    s = strip_accents(str(s)).upper()
    s = re.sub(r'[^A-Z\s]', ' ', s)
    return re.sub(r'\s+', ' ', s).strip()

NAME_DICT = set()

def set_name_dictionary(names):
    global NAME_DICT
    NAME_DICT = {norm(n) for n in names if isinstance(n, str) and n.strip()}

def tokens_match(a, b):
    """Two name tokens are equivalent if they're equal, if one is a single
    initial that matches the first letter of the other (handles abbreviations
    like 'B.' standing in for 'BRAGA'), or - for longer tokens only - if
    they're a near-identical spelling (handles a one-letter typo, like
    'OLIVERA' vs 'OLIVEIRA')."""
    if a == b:
        return True
    if len(a) == 1 and b.startswith(a):
        return True
    if len(b) == 1 and a.startswith(b):
        return True
    if len(a) >= 4 and len(b) >= 4 and SequenceMatcher(None, a, b).ratio() >= 0.82:
        return True
    return False

def name_tokens_match(ta, tb):
    """Compares two already-split, normalized name token lists by checking if
    the SHORTER one is an ordered subsequence of the LONGER one (each token
    matched via tokens_match, in the same relative order, gaps allowed
    anywhere - not just at the edges). This single rule covers every
    legitimate short-form pattern actually seen in these decks:
      - dropped/abbreviated MIDDLE names: 'AUGUSTO H. B. SILVA' vs
        'AUGUSTO HENRIQUE BRAGA SILVA'
      - a dropped TRAILING surname: 'LORRAYNE FERREIRA SILVA' vs
        'LORRAYNE FERREIRA SILVA ROSA MORAES'
      - going by a MIDDLE surname instead of the last one: 'EDUARDO CANOA'
        vs 'EDUARDO AGUIAR CANOA DE OLIVEIRA'
      - just first+last, dropping everything else: 'REGIS CARREIRA' vs
        'REGIS EDUARDO PRENHACA CARREIRA'
    ...while still rejecting two different people who share a first name
    and one coincidental token: 'JOAO PEDRO RODRIGUES' is NOT a subsequence
    of 'JOAO PAULO MORGADO FIZ RODRIGUES' because 'PEDRO' can't be found
    anywhere after 'JOAO' in that name.
    Returns (matched: bool, used_initial: bool, gaps: int) - gaps is how many
    tokens of the longer name had NO counterpart at all (fully dropped, like
    skipping 'BRASIL' entirely) - used to prefer the more complete match when
    two different real people both technically satisfy the subsequence rule."""
    if not ta or not tb:
        return False, False, 0
    shorter, longer = (ta, tb) if len(ta) <= len(tb) else (tb, ta)

    pos = 0
    used_initial = False
    for t in shorter:
        found = False
        for i in range(pos, len(longer)):
            if tokens_match(t, longer[i]):
                if t != longer[i]:
                    used_initial = True
                pos = i + 1
                found = True
                break
        if not found:
            return False, False, 0
    gaps = len(longer) - len(shorter)
    return True, used_initial, gaps

def line_matches_known_name(line_norm):
    """Returns (matched: bool, used_initial: bool)."""
    if not line_norm:
        return False, False
    if line_norm in NAME_DICT:
        return True, False
    tk = line_norm.split()
    first = tk[0]
    for n in NAME_DICT:
        n_tk = n.split()
        if not n_tk or n_tk[0] != first:
            continue
        matched, used_initial, _ = name_tokens_match(tk, n_tk)
        if matched:
            return True, used_initial
    return False, False

def clean(s):
    return re.sub(r'\s+', ' ', s).strip(" -|")

# ---------------------------------------------------------------------------
# Vínculo: rótulos usados dentro das caixas do PPT -> valor equivalente na
# coluna "Vínculo" do Fênix
# ---------------------------------------------------------------------------

VINCULO_LABELS = {
    'FUNCION': 'Funcionário',
    'ESTAGI': 'Estagiário',
    'ESTATUT': 'Estatutário',
    'PRESTADOR': 'Prestador de Serviço',
}

def label_for_line(line_upper):
    for key, val in VINCULO_LABELS.items():
        if line_upper.startswith(key):
            return val
    return None

NOISE_HEADERS = ["VISÃO DETALHADA", "ESTRUTURA MACRO", "ESTRUTURA DETALHADA", "VISÃO MACRO", "ORGANOGRAMA"]

def is_noise(text_upper):
    if text_upper.isdigit():
        return True
    if any(n in text_upper for n in NOISE_HEADERS):
        return True
    if re.match(r'^[A-ZÇÃÕÊÁÉÍÓÚ]+\s*[∙\.]\s*\d{4}$', text_upper):
        return True
    return False

# ---------------------------------------------------------------------------
# Texto alternativo: só reconhece o padrão "Gestor=NOME" - ignora qualquer
# outro texto alternativo (ex: descrição automática de imagem gerada por IA).
# ---------------------------------------------------------------------------

def get_gestor_alt_text(shape):
    cNvPr = shape._element.find(f'.//{NS_P}cNvPr')
    if cNvPr is None:
        return None
    descr = cNvPr.get('descr')
    if not descr:
        return None
    m = re.search(r'^Gestor\s*=\s*(.+)$', descr.strip(), re.IGNORECASE | re.MULTILINE)
    if not m:
        return None
    return re.sub(r'\s+', ' ', m.group(1)).strip()

def get_vinculo_alt_text(shape):
    """Lê 'Vinculo=X' (ou 'Vínculo=X') do texto alternativo - usado só para
    marcar o Estatutário do topo do organograma, sem depender de uma lista de
    nomes fixa no código."""
    cNvPr = shape._element.find(f'.//{NS_P}cNvPr')
    if cNvPr is None:
        return None
    descr = cNvPr.get('descr')
    if not descr:
        return None
    m = re.search(r'^V[ií]nculo\s*=\s*(.+)$', descr.strip(), re.IGNORECASE | re.MULTILINE)
    if not m:
        return None
    return re.sub(r'\s+', ' ', m.group(1)).strip()

def get_area_alt_text(shape):
    """Lê 'Area=X' (ou 'Área=X') do texto alternativo - ÚNICA fonte da área
    própria de uma caixa (cabeçalho de coordenador OU lista de funcionários).
    Não há mais adivinhação por texto visível: se a tag não existir, a caixa
    simplesmente não contribui nenhum nível de área. Pode conviver com
    'Gestor=...' na mesma caixa, cada tag na sua linha."""
    cNvPr = shape._element.find(f'.//{NS_P}cNvPr')
    if cNvPr is None:
        return None
    descr = cNvPr.get('descr')
    if not descr:
        return None
    m = re.search(r'^[AÁ]rea\s*=\s*(.+)$', descr.strip(), re.IGNORECASE | re.MULTILINE)
    if not m:
        return None
    return re.sub(r'\s+', ' ', m.group(1)).strip()

# ---------------------------------------------------------------------------
# Extração: percorre toda caixa de texto (AUTO_SHAPE) do PPT, sem geometria
# e sem hierarquia - só coleta trios (nome, vínculo, slide).
# ---------------------------------------------------------------------------

def extract_names_vinculo(pptx_path, known_names=None):
    if known_names is not None:
        set_name_dictionary(known_names)

    p = Presentation(pptx_path)
    rows = []

    # Pré-escaneamento: a tag "Vinculo=..." pode estar em só UMA das várias
    # ocorrências da mesma pessoa no deck (ex: só no slide de visão macro).
    # Escaneia o arquivo inteiro primeiro pra montar nome_norm -> vínculo,
    # e aplica esse valor em TODAS as ocorrências dessa pessoa depois.
    vinculo_overrides = {}

    def prescan(shapes):
        for shp in shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                prescan(shp.shapes)
                continue
            if shp.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX) or not getattr(shp, "has_text_frame", False):
                continue
            vinc_alt = get_vinculo_alt_text(shp)
            if not vinc_alt:
                continue
            txt = shp.text_frame.text
            txt_norm_breaks = txt.replace('\x0b', '\n').replace('\v', '\n')
            lines = [clean(l) for l in txt_norm_breaks.split('\n') if clean(l)]
            for l in lines:
                l_clean = re.sub(r'\(.*?\)', '', l).strip()
                matched, _ = line_matches_known_name(norm(l_clean)) if NAME_DICT else (False, False)
                if matched:
                    vinculo_overrides[norm(l_clean)] = vinc_alt
                    break

    for s in p.slides:
        prescan(s.shapes)

    def walk(shapes, slide_num):
        for shp in shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shp.shapes, slide_num)
                continue
            if shp.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
                continue
            if not getattr(shp, "has_text_frame", False):
                continue
            txt = shp.text_frame.text
            if not txt.strip():
                continue
            up = txt.strip().upper()
            if is_noise(up):
                continue

            gestor_alt = get_gestor_alt_text(shp)
            area_alt_roster = get_area_alt_text(shp)

            # PowerPoint's "quebra suave" (Shift+Enter) shows up as \x0b/\v inside
            # a paragraph, not as \n - normalize both to real line breaks first.
            txt_norm_breaks = txt.replace('\x0b', '\n').replace('\v', '\n')
            lines = [clean(l) for l in txt_norm_breaks.split('\n') if clean(l)]
            has_roster_label = any(label_for_line(l.upper()) for l in lines)

            if has_roster_label:
                # Caixa de lista: rótulos "FUNCIONÁRIOS:" / "ESTAGIÁRIOS:",
                # um nome por linha abaixo de cada rótulo.
                current_vinculo = None
                entries = []  # (vinculo, nome_bruto)
                for l in lines:
                    lab = label_for_line(l.upper())
                    if lab:
                        current_vinculo = lab
                        continue
                    if current_vinculo:
                        entries.append((current_vinculo, l))

                # Às vezes a maioria de uma lista tem uma "tag" de time grudada
                # no final do nome (ex: '... PO'), que não é sobrenome de
                # ninguém - mas nem sempre 100% das linhas têm a tag. Se a
                # palavra final mais comum aparece em pelo menos 60% das
                # linhas (com 3+ linhas no total), é ruído: tira só de quem
                # realmente termina com ela.
                if len(entries) >= 3:
                    from collections import Counter
                    last_words = [nome.split()[-1].upper() for _, nome in entries if len(nome.split()) > 1]
                    if last_words:
                        common_word, count = Counter(last_words).most_common(1)[0]
                        if count / len(entries) >= 0.6:
                            entries = [
                                (v, ' '.join(nome.split()[:-1])
                                 if len(nome.split()) > 1 and nome.split()[-1].upper() == common_word
                                 else nome)
                                for v, nome in entries
                            ]

                for vinc, nome in entries:
                    rows.append(dict(nome=nome, vinculo_pptx=vinc, gestor_pptx=gestor_alt,
                                      area_extra=area_alt_roster,
                                      slide=slide_num, tem_gestor_alt_text=gestor_alt is not None))
            else:
                # Caixa de cabeçalho: geralmente "NOME \n CARGO/ÁREA" (ordem varia).
                # Todo mundo em cabeçalho é Funcionário por padrão, exceto:
                # 1) quem tiver "Vinculo=..." no texto alternativo (fonte de
                #    verdade, quando presente); 2) senão, quem estiver na
                #    lista NOMES_ESTATUTARIO (fallback pra decks sem a tag).
                for l in lines:
                    l_clean = re.sub(r'\(.*?\)', '', l).strip()
                    matched, _ = line_matches_known_name(norm(l_clean)) if NAME_DICT else (False, False)
                    if matched:
                        override = vinculo_overrides.get(norm(l_clean))
                        if override:
                            vinc = override
                        else:
                            is_estatutario = any(
                                name_tokens_match(norm(l_clean).split(), norm(n).split())[0]
                                for n in NOMES_ESTATUTARIO
                            )
                            vinc = 'Estatutário' if is_estatutario else 'Funcionário'
                        rows.append(dict(nome=l_clean, vinculo_pptx=vinc, gestor_pptx=gestor_alt,
                                          slide=slide_num, tem_gestor_alt_text=gestor_alt is not None))

    for i, s in enumerate(p.slides):
        walk(s.shapes, i + 1)

    return pd.DataFrame(rows)


# ---------------------------------------------------------------------------
# Comparação com a base Fênix
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# Área: sobe a cadeia de GESTOR (já confiável, via texto alternativo) e pega a
# área própria de cada gestor no caminho, lida direto da caixa dele mesmo -
# sem depender de geometria/aninhamento (isso já foi tentado e não generalizou
# bem entre organogramas com layouts diferentes).
# ---------------------------------------------------------------------------

ROMAN_OR_NUM = {'I', 'II', 'III', 'IV', 'V', 'VI', 'VII', 'VIII', 'IX', 'X'}
STOPWORDS_AREA = {'PL', 'CORP'}

def split_suffix(tokens):
    """Separa um sufixo numérico/romano do final (ex: ['BELO','HORIZONTE','II']
    -> (['BELO','HORIZONTE'], ['II']))."""
    tokens = list(tokens)
    suffix = []
    while tokens and (tokens[-1] in ROMAN_OR_NUM or tokens[-1].isdigit()):
        suffix.insert(0, tokens.pop())
    return tokens, suffix

def area_token_matches(a, b):
    """Como tokens_match, mas TAMBÉM aceita truncamento de uma palavra
    (ex: 'RIB' por 'RIBEIRAO') - seguro aqui porque é só entre nomes de área,
    não nomes de pessoa (onde isso arriscaria confundir sobrenomes parecidos)."""
    if tokens_match(a, b):
        return True
    if len(a) >= 3 and len(b) >= 3:
        shorter, longer = (a, b) if len(a) <= len(b) else (b, a)
        if longer.startswith(shorter):
            return True
    return False

def area_tokens_match(ta, tb):
    """Mesma lógica de subsequência do name_tokens_match, mas usando
    area_token_matches (tolera truncamento de palavra)."""
    if not ta or not tb:
        return False
    shorter, longer = (ta, tb) if len(ta) <= len(tb) else (tb, ta)
    pos = 0
    for t in shorter:
        found = False
        for i in range(pos, len(longer)):
            if area_token_matches(t, longer[i]):
                pos = i + 1
                found = True
                break
        if not found:
            return False
    return True

def area_matches(pptx_area, fenix_value):
    """Compara um nível de área do PPT com um valor de coluna Abertura do
    Fênix, tolerando: abreviação normal, truncamento de palavra (ex: 'RIB.'
    = Ribeirão), e SIGLA de várias palavras (ex: 'BH' = Belo Horizonte),
    incluindo sufixo numérico/romano (ex: 'BH II' = '... Belo Horizonte II')."""
    pn_tokens = norm(pptx_area).split()
    fn_tokens = [t for t in norm(fenix_value).split() if t not in STOPWORDS_AREA]
    if area_tokens_match(pn_tokens, fn_tokens):
        return True
    pn_core, pn_suffix = split_suffix(pn_tokens)
    fn_core, fn_suffix = split_suffix(fn_tokens)
    if not pn_core or not fn_core:
        return False
    acronym = ''.join(w[0] for w in fn_core)
    if pn_core[0] != acronym or len(acronym) < 2:
        return False
    if len(pn_core) > 1 and not area_tokens_match(pn_core[1:], fn_core):
        return False
    if pn_suffix == fn_suffix:
        return True
    # "I" costuma ser o time "padrão/sem número" nesse tipo de deck, que o
    # Fênix às vezes só registra sem sufixo nenhum
    if pn_suffix == ['I'] and not fn_suffix:
        return True
    return False

def build_person_map(pptx_path, known_names):
    """Para CADA caixa de cabeçalho (nome+área) do deck inteiro: guarda a área
    própria (lida SÓ do texto alternativo 'Area=...', sem adivinhação por
    texto visível) e o gestor próprio (via 'Gestor=...'), por pessoa."""
    set_name_dictionary(known_names)
    p = Presentation(pptx_path)
    person_map = {}  # nome_norm -> {'nome':, 'area':, 'gestor':}

    def walk(shapes):
        for shp in shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk(shp.shapes)
                continue
            if shp.shape_type not in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX):
                continue
            if not getattr(shp, "has_text_frame", False):
                continue
            txt = shp.text_frame.text
            if not txt.strip() or is_noise(txt.strip().upper()):
                continue
            txt_norm_breaks = txt.replace('\x0b', '\n').replace('\v', '\n')
            lines = [clean(l) for l in txt_norm_breaks.split('\n') if clean(l)]
            if any(label_for_line(l.upper()) for l in lines):
                continue  # é lista de funcionarios, não cabeçalho

            nome_line = None
            for l in lines:
                l_clean = re.sub(r'\(.*?\)', '', l).strip()
                matched, _ = line_matches_known_name(norm(l_clean)) if NAME_DICT else (False, False)
                if matched and nome_line is None:
                    nome_line = l_clean
            if nome_line is None:
                continue

            gestor_alt = get_gestor_alt_text(shp)
            area_alt = get_area_alt_text(shp)
            key = norm(nome_line)
            existing = person_map.get(key, {})
            person_map[key] = dict(
                nome=nome_line,
                area=area_alt or existing.get('area'),
                gestor=gestor_alt or existing.get('gestor'),
            )

    for s in p.slides:
        walk(s.shapes)
    return person_map

def split_area_levels(texto):
    """Uma tag 'Area=X;Y;Z' pode conter mais de um nível numa caixa só - útil
    pra caixas 'banner' sem gestor próprio (ex: 'RESEARCH SALES GLOBAL' sem
    coordenador nomeado), que nunca apareceriam sozinhas na cadeia. Usa ';'
    como separador (não '/', que já aparece dentro de nomes de área reais,
    tipo 'CENTRO / NORTE')."""
    if not texto:
        return []
    return [p.strip() for p in texto.split(';') if p.strip()]

def build_area_chain(gestor_start, person_map, max_depth=10):
    """Sobe a cadeia de gestores (nome -> gestor -> gestor do gestor -> ...)
    juntando a área própria de cada um, do mais amplo (topo) pro mais
    específico. Sem geometria - só nome, indexado no person_map."""
    chain = []
    current = gestor_start
    seen = set()
    depth = 0
    while current and depth < max_depth:
        key = norm(current)
        if key in seen:
            break  # proteção contra ciclo
        seen.add(key)
        info = person_map.get(key)
        if info is None:
            tk = key.split()
            first = tk[0] if tk else ''
            for k, v in person_map.items():
                kt = k.split()
                if kt and kt[0] == first and name_tokens_match(tk, kt)[0]:
                    info = v
                    break
        if info is None:
            break
        if info.get('area'):
            for nivel in reversed(split_area_levels(info['area'])):
                chain.insert(0, nivel)
        current = info.get('gestor')
        depth += 1
    return chain


def find_match(nome_norm, pptx_dedup):
    """Returns (row_or_None, used_initial: bool, near_miss: bool).
    near_miss=True means the match only came from the fuzzy spelling fallback
    (e.g. 'CORREA' vs 'CORREIA') and should be treated as needing confirmation,
    never as an auto-confirmed OK.

    A person sometimes appears with slightly different spellings on different
    slides (e.g. 'SHEILA ANDRADE DE LIRA' on one slide, 'SHEILA DE LIRA' on
    another) - and only ONE of those spellings might carry the Gestor= data.
    So instead of returning the first match found (which could be an exact
    string match that happens to lack gestor_pptx), this collects every
    candidate (exact and fuzzy) and prefers whichever one actually has
    gestor_pptx filled."""
    tk = nome_norm.split()
    first = tk[0] if tk else ''
    candidates = []  # (row, used_initial, gaps)

    exact = pptx_dedup[pptx_dedup.nome_norm == nome_norm]
    for _, row in exact.iterrows():
        candidates.append((row, False, 0))

    for _, row in pptx_dedup.iterrows():
        if row['nome_norm'] == nome_norm:
            continue  # já coberto pelo exact acima
        pn_tokens = row['nome_norm'].split()
        if not pn_tokens or pn_tokens[0] != first:
            continue
        matched, used_initial, gaps = name_tokens_match(tk, pn_tokens)
        if matched:
            candidates.append((row, used_initial, gaps))

    if candidates:
        def tem_gestor(row):
            g = row.get('gestor_pptx')
            return not (g is None or (isinstance(g, float) and pd.isna(g)) or g == '')
        # prioriza: 1) tem gestor preenchido; 2) menos "buracos" (palavras do
        # nome completo que ficaram sem par nenhum - pular um sobrenome
        # inteiro tipo "BRASIL" é pior sinal do que só abreviar uma letra);
        # 3) não usou abreviação
        candidates.sort(key=lambda c: (
            not tem_gestor(c[0]),
            c[2],
            c[1],
        ))
        best_row, used_initial, _ = candidates[0]
        return best_row, used_initial, False

    # Nothing matched via tokens (first name itself may be misspelled, e.g.
    # "MARINA" vs "MARIANA", or a surname like "CORREA" vs "CORREIA"). Try a
    # pure character-similarity fallback and, if something is close enough,
    # surface it as a candidate to CONFIRM MANUALLY rather than trusting it.
    best_row, best_ratio = None, 0.0
    for _, row in pptx_dedup.iterrows():
        ratio = SequenceMatcher(None, nome_norm, row['nome_norm']).ratio()
        if ratio > best_ratio:
            best_ratio, best_row = ratio, row
    if best_row is not None and best_ratio >= 0.85:
        return best_row, False, True

    return None, False, False


def main():
    # 1. base Fênix
    fenix = pd.read_excel(CAMINHO_FENIX)
    fenix['nome_norm'] = fenix['Nome'].map(norm)
    fenix['gestor_norm'] = fenix['Gestor'].map(lambda x: norm(x) if pd.notna(x) else None)
    known_names = pd.concat([fenix['Nome'], fenix['Gestor']]).dropna().unique().tolist()

    # 2. extrai nome + vínculo + gestor (texto alternativo) do PPT
    pptx = extract_names_vinculo(CAMINHO_PPTX, known_names=known_names)
    pptx['nome_norm'] = pptx['nome'].map(norm)
    if 'area_extra' not in pptx.columns:
        pptx['area_extra'] = None

    # mapa de área própria por pessoa (pra montar a cadeia de área depois)
    person_map = build_person_map(CAMINHO_PPTX, known_names)

    print(f"Total de pessoas extraídas do PPT: {len(pptx)}")
    print(f"  - com Gestor= preenchido no texto alternativo: {pptx['tem_gestor_alt_text'].sum()}")
    print(f"  - sem Gestor= (não entram na validação de gestor): {(~pptx['tem_gestor_alt_text']).sum()}")
    print()

    # Dedupe: a mesma pessoa pode aparecer em mais de um slide (visão macro +
    # slide detalhado dela). Mantém o vínculo e o gestor quando houver mais de
    # um registro.
    pptx_dedup = (
        pptx.sort_values(['vinculo_pptx', 'tem_gestor_alt_text'], na_position='last', ascending=[True, False])
            .drop_duplicates(subset='nome_norm', keep='first')
            [['nome', 'nome_norm', 'vinculo_pptx', 'gestor_pptx', 'area_extra', 'slide']]
    )

    # 3. compara cada pessoa do Fênix com o que foi achado no PPT
    results = []
    for _, r in fenix.iterrows():
        m, used_initial, near_miss = find_match(r['nome_norm'], pptx_dedup)

        if m is None:
            status = 'NAO ENCONTRADO NO PPTX'
            vinc_pptx, gestor_pptx, slide = None, None, None
            status_gestor = 'N/A'
            area_pptx, status_area, status_area_rigoroso = None, 'N/A', 'N/A'
        else:
            vinc_pptx, gestor_pptx, slide = m['vinculo_pptx'], m['gestor_pptx'], m['slide']
            if near_miss:
                status = 'CONFERIR - POSSIVEL DIVERGENCIA DE GRAFIA'
            elif pd.isna(vinc_pptx):
                status = 'SEM VINCULO ROTULADO NO PPTX'
            elif vinc_pptx == r['Vínculo']:
                status = 'OK'
            else:
                status = 'DIVERGENTE'

            # Comparação de GESTOR - só faz sentido se a caixa tinha o texto
            # alternativo "Gestor=..." preenchido.
            if pd.isna(gestor_pptx) or not gestor_pptx:
                status_gestor = 'SEM GESTOR NO PPTX'
            elif pd.isna(r['gestor_norm']):
                status_gestor = 'SEM GESTOR NO FENIX'
            else:
                gestor_ok, _, _ = name_tokens_match(norm(gestor_pptx).split(), r['gestor_norm'].split())
                status_gestor = 'OK' if gestor_ok else 'DIVERGENTE'

            # Área: sobe a cadeia de gestores (já resolvida acima) juntando a
            # área própria de cada um; inclui também a área própria da
            # PESSOA, se ela mesma tiver cabeçalho (ex: ela é coordenadora).
            chain = build_area_chain(gestor_pptx, person_map) if gestor_pptx else []
            own_info = person_map.get(m['nome_norm'])
            if own_info and own_info.get('area'):
                chain = chain + split_area_levels(own_info['area'])
            # 'Area=' na própria caixa de FUNCIONÁRIOS/ESTAGIÁRIOS dela - o(s)
            # nível(is) mais específico(s), direto de onde ela está listada.
            # Também serve pra cobrir um nível "banner" sem gestor próprio
            # (ex: 'Area=RESEARCH SALES GLOBAL;RESEARCH SALES BR').
            area_extra = m.get('area_extra')
            if pd.notna(area_extra) and area_extra:
                chain = chain + split_area_levels(area_extra)
            chain = [c for c in chain if not c.upper().startswith('DIRETORIA')]
            # Remove repetição consecutiva (ex: o gestor já é da área "Centro
            # Oeste" e a própria pessoa também foi marcada "Centro Oeste" -
            # sem isso, o nível apareceria duplicado na cadeia).
            chain_dedup = []
            for c in chain:
                if not chain_dedup or norm(chain_dedup[-1]) != norm(c):
                    chain_dedup.append(c)
            chain = chain_dedup

            fenix_area_values = [str(r[c]) for c in ABERTURA_COLS if pd.notna(r[c]) and str(r[c]) != '-']
            # "Área Gerencial" é o nome geral do organograma inteiro (ex:
            # "Corporate") - nunca aparece escrito numa caixa específica, então
            # não entra na exigência da versão rigorosa, só no texto de
            # referência do relatório.
            fenix_check_values = [str(r[c]) for c in ABERTURA_COLS[1:] if pd.notna(r[c]) and str(r[c]) != '-']
            if not chain and len(fenix_check_values) == 0:
                status_area = 'OK'
                status_area_rigoroso = 'OK'
            elif not chain:
                status_area = 'SEM AREA NO PPTX'
                status_area_rigoroso = 'SEM AREA NO PPTX'
            else:
                # TOLERANTE: cada nível da cadeia do PPT precisa bater com
                # algum valor do Fênix - com 1 nível de tolerância, pra não
                # falhar por causa de um rótulo a mais que o Fênix não
                # rastreia (ex: "PRODUTOS PAGAMENTOS").
                niveis_batidos = sum(1 for level in chain if any(area_matches(level, fv) for fv in fenix_area_values))
                sem_par = len(chain) - niveis_batidos
                status_area = 'OK' if niveis_batidos >= 1 and sem_par <= 1 else 'AREA DIVERGENTE'
                # RIGOROSA: cada nível que o FÊNIX espera precisa aparecer em
                # algum lugar da cadeia do PPT - sem tolerância. Pega lacunas
                # de detalhe que a tolerante deixa passar (ex: organograma não
                # detalha até o nível de cidade/sub-time pra essa pessoa).
                niveis_ok = sum(1 for fv in fenix_check_values if any(area_matches(level, fv) for level in chain))
                status_area_rigoroso = 'OK' if niveis_ok == len(fenix_check_values) else 'AREA DIVERGENTE'
            area_pptx = ' / '.join(chain) if chain else None

        obs = 'Nome abreviado no PPT' if used_initial else ('Nome do PPT: ' + str(m['nome']) if (near_miss and m is not None) else '')
        results.append(dict(
            nome=r['Nome'], nome_pptx=m['nome'] if m is not None else None,
            vinculo_fenix=r['Vínculo'], vinculo_pptx=vinc_pptx, status_vinculo=status,
            gestor_fenix=r['Gestor'], gestor_pptx=gestor_pptx, status_gestor=status_gestor,
            abertura_fenix=' / '.join(str(r[c]) for c in ABERTURA_COLS if pd.notna(r[c]) and str(r[c]) != '-'),
            area_pptx=area_pptx, status_area=status_area, status_area_rigoroso=status_area_rigoroso,
            slide_pptx=slide, conferir_manualmente=obs
        ))

    res = pd.DataFrame(results)
    res.to_excel(CAMINHO_SAIDA, index=False)

    print("--- Vínculo ---")
    print(res.status_vinculo.value_counts())
    print()
    print("--- Gestor ---")
    print(res.status_gestor.value_counts())
    print()
    print("--- Área (TOLERANTE - até 1 nível sem correspondência) ---")
    print(res.status_area.value_counts())
    print()
    print("--- Área (RIGOROSA - exige todo nível do Fênix na cadeia) ---")
    print(res.status_area_rigoroso.value_counts())
    print()
    print("Divergências de gestor (agrupadas por par gestor_pptx -> gestor_fenix):")
    div = res[res.status_gestor == 'DIVERGENTE']
    if len(div):
        print(div.groupby(['gestor_pptx', 'gestor_fenix']).size().reset_index(name='qtd').to_string())
    print()
    print(f"Arquivo salvo em: {CAMINHO_SAIDA}")


if __name__ == "__main__":
    main()
